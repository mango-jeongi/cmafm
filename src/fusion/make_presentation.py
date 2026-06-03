import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(template_path=None, output_path="CS3315_Final_Presentation.pptx"):
    import os
    from dotenv import load_dotenv
    repo_root = Path(__file__).resolve().parents[2]
    load_dotenv(dotenv_path=repo_root / ".env")
    
    figures_dir_env = os.environ.get("CMAFM_FIGURES_DIR", "../figures")
    figures_dir_path = Path(figures_dir_env)
    if not figures_dir_path.is_absolute():
        figures_dir_path = (repo_root / figures_dir_path).resolve()
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Helpers
    def add_slide(title, layout_idx=1):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        return slide

    # Slide 1: Title
    slide = add_slide("Robust Object Detection in Low-Illumination via RGB-Thermal Sensor Fusion", layout_idx=0)
    subtitle = slide.placeholders[1]
    subtitle.text = "CS3315 Final Project: Introduction to Machine Learning and Big Data\nPresenter: Mingyu Jeong"

    # Slide 2: Introduction
    slide = add_slide("Introduction & Objective")
    tf = slide.placeholders[1].text_frame
    tf.text = "Objective: Enhance object detection reliability across diverse environmental conditions (night, fog, smoke)."
    tf.add_paragraph().text = "• Visible cameras (RGB) excel in daylight but fail in low illumination."
    tf.add_paragraph().text = "• Thermal cameras (LWIR) map heat signatures but lack texture."
    tf.add_paragraph().text = "• Goal: Effectively fuse these heterogeneous distributions under a unified detection framework."

    # Slide 3: Selection of Data
    slide = add_slide("Selection of Data")
    tf = slide.placeholders[1].text_frame
    tf.text = "Datasets Used:"
    tf.add_paragraph().text = "• M3FD: 4,200 aligned RGB-Thermal pairs (diverse weather)."
    tf.add_paragraph().text = "• FLIR ADAS v2: 5,142 aligned pairs (urban driving)."
    tf.add_paragraph().text = "Data Munging & Cleaning:"
    tf.add_paragraph().text = "• Class matching across 6 standard categories."
    tf.add_paragraph().text = "• Spatial verification to ensure pixel-perfect overlap."

    # Slide 4: Feature Engineering (Day/Night Split)
    slide = add_slide("Feature Engineering: Illumination Splitting")
    tf = slide.placeholders[1].text_frame
    tf.text = "To rigorously test low-illumination robustness, we engineered an automated brightness splitter:"
    tf.add_paragraph().text = "• Calculated average pixel intensity of the visible (RGB) channel."
    tf.add_paragraph().text = "• Night Split: Grayscale intensity < 60."
    tf.add_paragraph().text = "• Day Split: Grayscale intensity >= 60."
    tf.add_paragraph().text = "This allows us to evaluate dynamic context sensitivity."

    # Slide 5: Methods, APIs, and Tools
    slide = add_slide("Methods, APIs, & Tools")
    tf = slide.placeholders[1].text_frame
    tf.text = "Core Framework: PyTorch v2.4.0 & torchvision"
    tf.add_paragraph().text = "• torchvision FasterRCNN: Provided a customizable backbone hook for modular design."
    tf.add_paragraph().text = "• Albumentations: High-performance data augmentation."
    tf.add_paragraph().text = "• pycocotools: Standardized mAP calculations."
    tf.add_paragraph().text = "• SLURM: Job orchestration on Hamming HPC."

    # Slide 6: Model Architecture
    slide = add_slide("Model Architecture: CMAFM")
    tf = slide.placeholders[1].text_frame
    tf.text = "Cross-Modal Attention Fusion Module (CMAFM):"
    tf.add_paragraph().text = "1. Channel Cross-Attention: Compresses features via Global Average Pooling (GAP) to exchange global context."
    tf.add_paragraph().text = "2. Spatial Cross-Gating: Preserves high-frequency textures using depthwise separable convolutions."
    tf.add_paragraph().text = "3. Gated Fusion: Dynamically merges spatial maps using a location-wise learnable gate."
    
    img_path = figures_dir_path / "pipeline_cmafm.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(3.5), height=Inches(3.5))

    # Slide 7: Results (Ablation Study)
    slide = add_slide("Results: Ablation & Modality Comparison")
    tf = slide.placeholders[1].text_frame
    tf.text = "Performance on M3FD (Faster R-CNN):"
    tf.add_paragraph().text = "• Unimodal Baselines: RGB (62.3%) vs Thermal (53.4%)"
    tf.add_paragraph().text = "• Early Fusion (6-channel input): 65.3%"
    tf.add_paragraph().text = "• Late Fusion (Batched NMS): 65.7% (but slow: 42 FPS)"
    tf.add_paragraph().text = "• Full CMAFM (Ours): 71.7% mAP@0.5"

    # Slide 8: Results (Cross-Architecture SOTA)
    slide = add_slide("Results: Cross-Architecture SOTA")
    tf = slide.placeholders[1].text_frame
    tf.text = "Generalization to YOLOv5l (NVIDIA L40S):"
    tf.add_paragraph().text = "• YOLOv5l + CMAFM achieved an average of 85.75% mAP@0.5."
    tf.add_paragraph().text = "• Retained real-time inference latency (58 FPS)."
    tf.add_paragraph().text = "• Substantial gains (+6.55%) over Early Fusion."

    # Slide 9: Qualitative Visualizations
    slide = add_slide("Qualitative Visualizations")
    tf = slide.placeholders[1].text_frame
    tf.text = "Robust localization regardless of illumination."
    
    img_day = figures_dir_path / "fig_day_00872_clean.png"
    img_night = figures_dir_path / "fig_cmafm_night_clean.png"
    if img_day.exists() and img_night.exists():
        try:
            slide.shapes.add_picture(str(img_day), Inches(0.5), Inches(2.5), width=Inches(4))
            slide.shapes.add_picture(str(img_night), Inches(5), Inches(2.5), width=Inches(4))
        except Exception as e:
            print(f"Warning: Could not add images: {e}")

    # Slide 10: Discussion & Summary
    slide = add_slide("Discussion & Summary")
    tf = slide.placeholders[1].text_frame
    tf.text = "Key Findings:"
    tf.add_paragraph().text = "• Standard fusion strategies offer poor trade-offs between speed and performance."
    tf.add_paragraph().text = "• Global Average Pooling provides a highly parameter-efficient means to context-weight features."
    tf.add_paragraph().text = "• The PyTorch API allowed for rapid, modular experimentation of complex attention mechanisms."

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", type=str, default=None, help="Path to base PPTX template")
    args = parser.parse_args()
    create_presentation(args.template)
